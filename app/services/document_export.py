"""Document export (chat-module A.4). Turns an assistant answer (markdown) into a
downloadable file in the requested format:

    Markdown · Text · CSV · PDF · Word (.docx) · Excel (.xlsx) ·
    PowerPoint (.pptx) · zip (code-block archive)

``export_document(content, fmt, title) -> (bytes, mime_type, filename)``. A small
block parser extracts headings / paragraphs / lists / tables / code fences so each
generator can lay the content out sensibly.
"""
from __future__ import annotations

import csv
import io
import re
import zipfile
from typing import Optional

_FENCE_RE = re.compile(r"^```(\w*)\s*$")
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_LIST_RE = re.compile(r"^\s*([-*+]|\d+[.)])\s+")
_SEP_RE = re.compile(r"^\s*\|?[\s:|-]+\|?\s*$")
_QUOTE_RE = re.compile(r"^\s*>\s?(.*)$")           # blockquote → callout box

_LANG_EXT = {
    "python": ".py", "py": ".py", "javascript": ".js", "js": ".js",
    "typescript": ".ts", "ts": ".ts", "java": ".java", "c": ".c", "cpp": ".cpp",
    "c++": ".cpp", "csharp": ".cs", "cs": ".cs", "go": ".go", "rust": ".rs",
    "ruby": ".rb", "php": ".php", "swift": ".swift", "kotlin": ".kt",
    "dart": ".dart", "sql": ".sql", "bash": ".sh", "sh": ".sh", "shell": ".sh",
    "html": ".html", "css": ".css", "json": ".json", "yaml": ".yaml",
    "yml": ".yml", "xml": ".xml", "markdown": ".md", "md": ".md",
}


# ─── Lightweight markdown block parser ──────────────────────────────────────

def _parse_blocks(md: str) -> list[dict]:
    lines = (md or "").split("\n")
    n = len(lines)
    blocks: list[dict] = []
    i = 0
    while i < n:
        line = lines[i]
        fence = _FENCE_RE.match(line.strip())
        if fence:
            lang = fence.group(1)
            j, code = i + 1, []
            while j < n and not lines[j].strip().startswith("```"):
                code.append(lines[j])
                j += 1
            blocks.append({"type": "code", "lang": lang, "text": "\n".join(code)})
            i = j + 1
            continue
        # table: a "| … |" line followed by a separator row
        if "|" in line and i + 1 < n and _SEP_RE.match(lines[i + 1]) and "-" in lines[i + 1]:
            k, tbl = i, []
            while k < n and "|" in lines[k]:
                tbl.append(lines[k])
                k += 1
            rows = _parse_table(tbl)
            if rows:
                blocks.append({"type": "table", "rows": rows})
                i = k
                continue
        hm = _HEADING_RE.match(line)
        if hm:
            blocks.append({"type": "heading", "level": len(hm.group(1)), "text": hm.group(2).strip()})
            i += 1
            continue
        if _LIST_RE.match(line):
            items = []
            while i < n and _LIST_RE.match(lines[i]):
                items.append(_LIST_RE.sub("", lines[i]).strip())
                i += 1
            blocks.append({"type": "list", "items": items})
            continue
        if _QUOTE_RE.match(line):
            qlines = []
            while i < n and _QUOTE_RE.match(lines[i]):
                qlines.append(_QUOTE_RE.match(lines[i]).group(1))
                i += 1
            blocks.append({"type": "quote", "text": "\n".join(qlines).strip()})
            continue
        if not line.strip():
            i += 1
            continue
        para = []
        while (i < n and lines[i].strip()
               and not _FENCE_RE.match(lines[i].strip())
               and not _HEADING_RE.match(lines[i])
               and not _LIST_RE.match(lines[i])):
            para.append(lines[i].strip())
            i += 1
        # Keep the source line breaks (single newlines) rather than collapsing
        # them into one run-on line — many answers put structured content (e.g.
        # "Label: values" rows) on their own lines and rely on that layout.
        blocks.append({"type": "para", "text": "\n".join(para)})
    return blocks


def _parse_table(tbl_lines: list[str]) -> list[list[str]]:
    rows = []
    for l in tbl_lines:
        if _SEP_RE.match(l) and "-" in l:
            continue
        rows.append([c.strip() for c in l.strip().strip("|").split("|")])
    return rows


def _strip_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1 (\2)", text)
    return text.strip()


def _docx_add_runs(paragraph, text: str) -> None:
    """Append text to a docx paragraph, rendering inline **bold** as bold runs
    (other markers are stripped). Keeps emphasis instead of flattening it."""
    for seg in re.split(r"(\*\*.+?\*\*)", text):
        if not seg:
            continue
        if len(seg) > 4 and seg.startswith("**") and seg.endswith("**"):
            paragraph.add_run(_strip_md(seg[2:-2])).bold = True
        else:
            paragraph.add_run(_strip_md(seg))


def _norm(s: str) -> str:
    """Normalize for comparison (lowercase alphanumerics only)."""
    return re.sub(r"[^a-z0-9]+", "", _strip_md(s or "").lower())


def _skip_title(blocks: list[dict], title: str) -> bool:
    """True when the content already opens with the title heading, so the
    generator shouldn't add a second (duplicate) title on top."""
    return bool(blocks and blocks[0]["type"] == "heading"
                and _norm(blocks[0]["text"]) == _norm(title))


def _safe_title(title: str) -> str:
    title = re.sub(r"[^\w\-. ]+", "", (title or "document")).strip() or "document"
    return title[:60]


# ─── Strip the chat wrapper so the file is JUST the document ─────────────────
# Answers arrive wrapped in chat affordances — a "download with the Export
# button" note, a trailing "Follow-ups" section, and conversational preamble /
# a title that names the file format. None of that belongs in the downloaded
# file. Every model wraps differently, so we normalize on export.

_FOLLOWUP_HEAD_RE = re.compile(
    r"^\s*#{0,6}\s*\*{0,2}\s*follow[\s\-]?ups?\s*\*{0,2}\s*:?\s*$", re.IGNORECASE)
_FORMAT_PAREN_RE = re.compile(
    r"\((?:markdown|md|word|docx|pdf|excel|xlsx|csv|text|txt|powerpoint|pptx)\)\s*$",
    re.IGNORECASE)
_FRAMING_RE = re.compile(
    r"^\s*(here(?:'?s| is| you go)|sure[,!.: ]|certainly[,!.: ]|of course[,!.: ]|"
    r"absolutely[,!.: ]|below is\b|i(?:'?ve| have) (?:created|prepared|revised|"
    r"updated|drafted|put together))",
    re.IGNORECASE)
_SEP_ONLY_RE = re.compile(r"^\s*([-*_])\1{2,}\s*$")  # --- *** ___


def _is_download_note(line: str) -> bool:
    low = line.lower()
    if ("beneath this message" in low or "beneath your message" in low
            or "below this message" in low):
        return True
    return ("button" in low) and ("download" in low or "export" in low)


def _clean_for_export(md: str) -> str:
    """Drop the assistant's chat wrapper (download-button note, Follow-ups
    section, conversational preamble / format-title) so only the document
    remains. Falls back to the original if cleaning would empty it."""
    lines = (md or "").split("\n")

    # Cut a trailing "Follow-ups" section (its heading → end of message).
    for i, ln in enumerate(lines):
        if _FOLLOWUP_HEAD_RE.match(ln):
            lines = lines[:i]
            break

    # Remove any "download via the Export button" affordance line(s).
    lines = [ln for ln in lines if not _is_download_note(ln)]

    # Trim leading blanks / separators / conversational framing / format titles.
    while lines:
        s = lines[0].strip()
        bare = re.sub(r"^#{1,6}\s*", "", s).replace("*", "").strip()
        if (not s or _SEP_ONLY_RE.match(s) or _FRAMING_RE.match(s)
                or _FORMAT_PAREN_RE.search(bare)):
            lines.pop(0)
        else:
            break

    # Trim trailing blanks / separators.
    while lines and (not lines[-1].strip() or _SEP_ONLY_RE.match(lines[-1].strip())):
        lines.pop()

    cleaned = "\n".join(lines).strip()
    return cleaned or (md or "").strip()


def _clean_title(title: str) -> str:
    """Drop a trailing format tag like '(Markdown)' from a title/filename."""
    return _FORMAT_PAREN_RE.sub("", title or "").strip() or (title or "")


# The PDF core fonts are latin-1 only, so Unicode punctuation would otherwise
# become "?" — transliterate the common ones (bullets, dashes, smart quotes,
# non-breaking / thin spaces, arrows) to clean ASCII first.
_TRANSLIT = {
    0x2022: "-", 0x2023: "-", 0x25AA: "-", 0x25CF: "-", 0x2043: "-", 0x00B7: "-",
    0x2014: "-", 0x2013: "-", 0x2011: "-", 0x2212: "-",       # dashes / minus
    0x2018: "'", 0x2019: "'", 0x201A: "'", 0x2032: "'",       # single quotes
    0x201C: '"', 0x201D: '"', 0x201E: '"', 0x2033: '"',       # double quotes
    0x2026: "...", 0x2192: "->", 0x2190: "<-", 0x21D2: "=>", 0x21D0: "<=",
    0x00A0: " ", 0x2009: " ", 0x202F: " ", 0x2007: " ", 0x200B: "",  # spaces
    0x2013: "-", 0x2015: "-", 0xFE0F: "",
    # math / comparison operators (would otherwise show as "?")
    0x2264: "<=", 0x2265: ">=", 0x2260: "!=", 0x2248: "~=", 0x2261: "==",
    0x00D7: "x", 0x00F7: "/", 0x00B1: "+/-", 0x221E: "inf", 0x221A: "sqrt",
    0x00B0: " deg", 0x03BC: "u", 0x2211: "sum", 0x220F: "prod", 0x2208: "in",
    0x2200: "for all", 0x2203: "exists", 0x2205: "empty", 0x2229: "and",
    0x222A: "or", 0x00BD: "1/2", 0x00BC: "1/4", 0x00BE: "3/4", 0x2153: "1/3",
    0x2154: "2/3", 0x2013: "-", 0x2032: "'", 0x2033: '"',
}


def _latin(s: str) -> str:
    return (s or "").translate(_TRANSLIT).encode("latin-1", "replace").decode("latin-1")


# ─── Format generators ──────────────────────────────────────────────────────

def _to_markdown(content: str, title: str):
    return content.encode("utf-8"), "text/markdown", f"{title}.md"


def _to_text(content: str, title: str):
    out = []
    for b in _parse_blocks(content):
        t = b["type"]
        if t == "heading":
            out += [_strip_md(b["text"]).upper(), ""]
        elif t == "para":
            out += [_strip_md(b["text"]), ""]
        elif t == "list":
            out += ["  • " + _strip_md(x) for x in b["items"]] + [""]
        elif t == "quote":
            out += ["  | " + _strip_md(ln)
                    for ln in b["text"].split("\n") if ln.strip()] + [""]
        elif t == "code":
            out += [b["text"], ""]
        elif t == "table":
            out += ["   ".join(_strip_md(c) for c in r) for r in b["rows"]] + [""]
    return ("\n".join(out).strip() or content).encode("utf-8"), "text/plain", f"{title}.txt"


def _to_csv(content: str, title: str):
    blocks = _parse_blocks(content)
    table = next((b["rows"] for b in blocks if b["type"] == "table"), None)
    buf = io.StringIO()
    w = csv.writer(buf)
    if table:
        for r in table:
            w.writerow([_strip_md(c) for c in r])
    else:
        for b in blocks:
            if b["type"] == "para":
                for ln in b["text"].split("\n"):
                    if ln.strip():
                        w.writerow([_strip_md(ln)])
            elif b["type"] == "list":
                for x in b["items"]:
                    w.writerow([_strip_md(x)])
            elif b["type"] == "heading":
                w.writerow([_strip_md(b["text"])])
    return buf.getvalue().encode("utf-8"), "text/csv", f"{title}.csv"


# Document palette (kept subtle + professional).
_ACCENT_HEX = "1F4E79"        # dark blue — title, headings, table header
_CODE_BG = "F6F8FA"           # light gray — code box
_CODE_BORDER = "E1E4E8"
_CALLOUT_BG = "EAF1FB"        # light blue — callout box
_CALLOUT_BORDER = "C6D6EF"
_TABLE_BORDER = "BFBFBF"


def _docx_shade(cell, fill: str) -> None:
    """Set a table cell's background fill (python-docx has no direct API)."""
    try:
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), fill)
        tcPr.append(shd)
    except Exception:
        pass


def _docx_borders(table, color: str = _TABLE_BORDER, sz: str = "4") -> None:
    try:
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        tblPr = table._tbl.tblPr
        borders = OxmlElement("w:tblBorders")
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            el = OxmlElement(f"w:{edge}")
            el.set(qn("w:val"), "single")
            el.set(qn("w:sz"), sz)
            el.set(qn("w:color"), color)
            borders.append(el)
        tblPr.append(borders)
    except Exception:
        pass


def _docx_page_number(section) -> None:
    """Centered 'Page N' field in the footer."""
    try:
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        p = section.footer.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        lead = p.add_run("Page ")
        lead.font.size = Pt(8)
        lead.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        r = p.add_run()
        r.font.size = Pt(8)
        r.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        begin = OxmlElement("w:fldChar")
        begin.set(qn("w:fldCharType"), "begin")
        instr = OxmlElement("w:instrText")
        instr.set(qn("xml:space"), "preserve")
        instr.text = " PAGE "
        end = OxmlElement("w:fldChar")
        end.set(qn("w:fldCharType"), "end")
        r._r.append(begin)
        r._r.append(instr)
        r._r.append(end)
    except Exception:
        pass


def _docx_code(d, code: str) -> None:
    """Code in a shaded, bordered, monospace box."""
    from docx.shared import Pt, RGBColor
    table = d.add_table(rows=1, cols=1)
    table.autofit = True
    _docx_borders(table, color=_CODE_BORDER, sz="4")
    cell = table.rows[0].cells[0]
    _docx_shade(cell, _CODE_BG)
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(0)
    for idx, ln in enumerate(code.split("\n")):
        if idx:
            p.add_run().add_break()
        run = p.add_run(ln.rstrip())
        run.font.name = "Consolas"
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor(0x24, 0x29, 0x2E)
    d.add_paragraph().paragraph_format.space_after = Pt(2)


def _docx_callout(d, text: str) -> None:
    """A shaded, bordered callout box (e.g. Note / Important / Optimization)."""
    table = d.add_table(rows=1, cols=1)
    table.autofit = True
    _docx_borders(table, color=_CALLOUT_BORDER, sz="4")
    cell = table.rows[0].cells[0]
    _docx_shade(cell, _CALLOUT_BG)
    p = cell.paragraphs[0]
    for idx, ln in enumerate(l for l in text.split("\n") if l.strip()):
        if idx:
            p.add_run().add_break()
        _docx_add_runs(p, ln)
    d.add_paragraph()


def _docx_table(d, rows: list) -> None:
    """Table with a shaded, bold, white header row + light grid borders."""
    from docx.shared import RGBColor, Pt
    cols = max(len(r) for r in rows)
    table = d.add_table(rows=0, cols=cols)
    table.autofit = True
    _docx_borders(table)
    for ri, r in enumerate(rows):
        cells = table.add_row().cells
        for ci in range(cols):
            cell = cells[ci]
            cell.paragraphs[0].text = ""
            _docx_add_runs(cell.paragraphs[0], r[ci] if ci < len(r) else "")
            if ri == 0:
                _docx_shade(cell, _ACCENT_HEX)
                for run in cell.paragraphs[0].runs:
                    run.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.size = Pt(10)
    d.add_paragraph().paragraph_format.space_after = Pt(2)


def _docx_title(d, title: str, subtitle: Optional[str]) -> None:
    from docx.shared import Pt, RGBColor
    p = d.add_paragraph()
    r = p.add_run(_strip_md(title))
    r.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    if subtitle:
        sp = d.add_paragraph()
        sr = sp.add_run(_strip_md(subtitle))
        sr.font.size = Pt(13)
        sr.italic = True
        sr.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    d.add_paragraph()


def _to_docx(content: str, title: str):
    import docx
    from docx.shared import Pt, RGBColor, Inches
    d = docx.Document()

    # Base look: Calibri body, comfortable margins, page-number footer.
    normal = d.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    for sec in d.sections:
        sec.top_margin = Inches(0.9)
        sec.bottom_margin = Inches(0.9)
        sec.left_margin = Inches(1.0)
        sec.right_margin = Inches(1.0)
        _docx_page_number(sec)

    blocks = _parse_blocks(content)
    start = 0
    # Styled title: use the passed title, or promote the content's own first
    # heading (with a following short line as the subtitle).
    if not _skip_title(blocks, title):
        _docx_title(d, title, None)
    elif blocks:
        subtitle = None
        if (len(blocks) > 1 and blocks[1]["type"] == "para"
                and len(blocks[1]["text"]) <= 140 and "\n" not in blocks[1]["text"]):
            subtitle = blocks[1]["text"]
            start = 2
        else:
            start = 1
        _docx_title(d, _strip_md(blocks[0]["text"]), subtitle)

    for b in blocks[start:]:
        t = b["type"]
        if t == "heading":
            h = d.add_heading(level=min(b["level"], 4))
            _docx_add_runs(h, b["text"])
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        elif t == "para":
            plines = [ln for ln in b["text"].split("\n") if ln.strip()]
            if plines:
                p = d.add_paragraph()
                for idx, ln in enumerate(plines):
                    if idx:
                        p.add_run().add_break()
                    _docx_add_runs(p, ln)
        elif t == "list":
            for x in b["items"]:
                _docx_add_runs(d.add_paragraph(style="List Bullet"), x)
        elif t == "quote":
            _docx_callout(d, b["text"])
        elif t == "code":
            _docx_code(d, b["text"])
        elif t == "table" and b["rows"]:
            _docx_table(d, b["rows"])
    buf = io.BytesIO()
    d.save(buf)
    return (buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{title}.docx")


def _pdf_write(pdf, h: float, text: str, markdown: bool = False) -> None:
    """Write a block to the PDF, one source line per line (structure preserved),
    hard-wrapping very long unbreakable tokens so fpdf can always place them.
    With markdown=True, inline **bold** / __italic__ render as emphasis.

    new_x/new_y are explicit: each line must return to the left margin and move
    down, otherwise a second w=0 multi_cell starts at the right edge and fpdf
    raises "not enough horizontal space" (which would silently drop the line)."""
    from fpdf.enums import XPos, YPos
    for raw in (text or "").split("\n"):
        line = _latin(raw)
        # Insert a break opportunity inside tokens longer than ~90 chars.
        line = re.sub(r"(\S{90})", r"\1 ", line)
        try:
            pdf.multi_cell(0, h, line if line else " ", markdown=markdown,
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        except Exception:
            # Retry without markdown if a stray '*'/'_' trips the parser.
            try:
                pdf.multi_cell(0, h, line if line else " ",
                               new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            except Exception:
                continue


def _pdf_table(pdf, rows: list) -> None:
    """A bordered table with a shaded, white, bold header row. Falls back to
    pipe-joined lines if fpdf can't fit the layout (e.g. a very wide table)."""
    ncols = max((len(r) for r in rows), default=0)
    if ncols == 0:
        return
    data = [[_latin(_strip_md(c)) for c in r] + [""] * (ncols - len(r))
            for r in rows]
    try:
        from fpdf.fonts import FontFace
        headings = FontFace(emphasis="BOLD", color=(255, 255, 255),
                            fill_color=(31, 78, 121))
        with pdf.table(first_row_as_headings=True, headings_style=headings,
                       text_align="LEFT", line_height=6) as table:
            for r in data:
                trow = table.row()
                for cell in r:
                    trow.cell(cell)
    except Exception:
        pdf.set_font("Helvetica", "", 9)
        for r in rows:
            _pdf_write(pdf, 5, " | ".join(_strip_md(c) for c in r))


def _pdf_code(pdf, code: str) -> None:
    """Code in a shaded (light-gray) monospace block."""
    from fpdf.enums import XPos, YPos
    pdf.ln(1)
    pdf.set_font("Courier", "", 8)
    pdf.set_fill_color(244, 245, 247)
    pdf.set_text_color(36, 41, 46)
    for raw in (code or "").split("\n"):
        line = re.sub(r"(\S{100})", r"\1 ", _latin(raw.rstrip()))
        try:
            pdf.multi_cell(0, 4.6, line if line else " ", fill=True,
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        except Exception:
            continue
    pdf.set_text_color(0, 0, 0)
    pdf.ln(2)


def _pdf_callout(pdf, text: str) -> None:
    """A shaded (light-blue) callout box (Note / Important / Optimization…)."""
    from fpdf.enums import XPos, YPos
    pdf.ln(1)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_fill_color(234, 241, 251)
    pdf.set_text_color(20, 40, 70)
    for raw in (l for l in text.split("\n") if l.strip()):
        line = re.sub(r"(\S{90})", r"\1 ", _latin(raw))
        try:
            pdf.multi_cell(0, 6, line, markdown=True, fill=True,
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        except Exception:
            try:
                pdf.multi_cell(0, 6, line, fill=True,
                               new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            except Exception:
                continue
    pdf.set_text_color(0, 0, 0)
    pdf.ln(2)


def _to_pdf(content: str, title: str):
    from fpdf import FPDF

    class _StyledPDF(FPDF):
        def footer(self):  # centered 'Page N' on every page
            self.set_y(-12)
            self.set_font("Helvetica", "", 8)
            self.set_text_color(128, 128, 128)
            try:
                self.cell(0, 8, f"Page {self.page_no()}", align="C")
            except Exception:
                pass
            self.set_text_color(0, 0, 0)

    pdf = _StyledPDF()
    pdf.set_auto_page_break(True, margin=16)
    pdf.add_page()
    blocks = _parse_blocks(content)
    start = 0

    def _title(t: str, sub: Optional[str]) -> None:
        pdf.set_font("Helvetica", "B", 22)
        pdf.set_text_color(31, 78, 121)
        _pdf_write(pdf, 10, _strip_md(t))
        if sub:
            pdf.set_font("Helvetica", "I", 12)
            pdf.set_text_color(89, 89, 89)
            _pdf_write(pdf, 7, _strip_md(sub))
        pdf.set_text_color(0, 0, 0)
        pdf.ln(3)

    if not _skip_title(blocks, title):
        _title(title, None)
    elif blocks:
        sub = None
        if (len(blocks) > 1 and blocks[1]["type"] == "para"
                and len(blocks[1]["text"]) <= 140 and "\n" not in blocks[1]["text"]):
            sub = blocks[1]["text"]
            start = 2
        else:
            start = 1
        _title(_strip_md(blocks[0]["text"]), sub)

    for b in blocks[start:]:
        t = b["type"]
        if t == "heading":
            pdf.set_font("Helvetica", "B", max(11, 17 - b["level"]))
            pdf.set_text_color(31, 78, 121)
            _pdf_write(pdf, 7, _strip_md(b["text"]))
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif t == "para":
            pdf.set_font("Helvetica", "", 11)
            _pdf_write(pdf, 6, b["text"], markdown=True)
            pdf.ln(1)
        elif t == "list":
            pdf.set_font("Helvetica", "", 11)
            for x in b["items"]:
                _pdf_write(pdf, 6, "-  " + x, markdown=True)
            pdf.ln(1)
        elif t == "quote":
            _pdf_callout(pdf, b["text"])
        elif t == "code":
            _pdf_code(pdf, b["text"])
        elif t == "table" and b["rows"]:
            pdf.set_font("Helvetica", "", 9)
            _pdf_table(pdf, b["rows"])
            pdf.ln(2)
    return bytes(pdf.output()), "application/pdf", f"{title}.pdf"


def _to_xlsx(content: str, title: str):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    blocks = _parse_blocks(content)
    tables = [b["rows"] for b in blocks if b["type"] == "table"]
    if tables:
        for ti, rows in enumerate(tables):
            sheet = ws if ti == 0 else wb.create_sheet(f"Table {ti + 1}")
            for r in rows:
                sheet.append([_strip_md(c) for c in r])
    else:
        for b in blocks:
            if b["type"] == "heading":
                ws.append([_strip_md(b["text"])])
            elif b["type"] == "para":
                for ln in b["text"].split("\n"):
                    if ln.strip():
                        ws.append([_strip_md(ln)])
            elif b["type"] == "list":
                for x in b["items"]:
                    ws.append([_strip_md(x)])
    buf = io.BytesIO()
    wb.save(buf)
    return (buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            f"{title}.xlsx")


def _to_pptx(content: str, title: str):
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = title

    state = {"title": "Overview", "bullets": []}

    def flush():
        if not state["bullets"]:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = state["title"][:80]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bi, bl in enumerate(state["bullets"][:10]):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            p.text = bl[:200]
        state["bullets"] = []

    for b in _parse_blocks(content):
        if b["type"] == "heading":
            flush()
            state["title"] = _strip_md(b["text"])
        elif b["type"] in ("para", "quote"):
            state["bullets"].extend(
                _strip_md(ln) for ln in b["text"].split("\n") if ln.strip())
        elif b["type"] == "list":
            state["bullets"].extend(_strip_md(x) for x in b["items"])
    flush()
    buf = io.BytesIO()
    prs.save(buf)
    return (buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            f"{title}.pptx")


def _to_zip(content: str, title: str):
    blocks = _parse_blocks(content)
    codes = [b for b in blocks if b["type"] == "code" and b["text"].strip()]
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("README.md", content or "")
        seen: dict[str, int] = {}
        for b in codes:
            ext = _LANG_EXT.get((b["lang"] or "").lower(), ".txt")
            base = (b["lang"] or "file").lower()
            seen[base] = seen.get(base, 0) + 1
            z.writestr(f"project/{base}_{seen[base]}{ext}", b["text"])
    return buf.getvalue(), "application/zip", f"{title}.zip"


_GENERATORS = {
    "md": _to_markdown, "markdown": _to_markdown,
    "txt": _to_text, "text": _to_text,
    "csv": _to_csv,
    "pdf": _to_pdf,
    "docx": _to_docx, "word": _to_docx,
    "xlsx": _to_xlsx, "excel": _to_xlsx,
    "pptx": _to_pptx, "powerpoint": _to_pptx,
    "zip": _to_zip,
}

SUPPORTED_FORMATS = ["markdown", "word", "pdf", "excel", "csv", "text", "powerpoint", "zip"]


def export_document(content: str, fmt: str, title: str = "document",
                    clean: bool = True):
    """Return (bytes, mime_type, filename) for the requested format.

    clean=True (the in-response download) strips the chat wrapper so the file
    holds only the requested document. clean=False (the Export menu) keeps the
    whole AI response verbatim.
    """
    gen = _GENERATORS.get((fmt or "md").lower().strip())
    if gen is None:
        raise ValueError(f"Unsupported export format: {fmt}")
    body = _clean_for_export(content or "") if clean else (content or "")
    return gen(body, _safe_title(_clean_title(title)))
